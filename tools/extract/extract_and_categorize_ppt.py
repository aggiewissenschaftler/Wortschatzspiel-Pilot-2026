from __future__ import annotations

import csv
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Dict, Iterable, List, Tuple

from pptx import Presentation
import yaml

# -----------------------------
# Lexicon candidate filter
# -----------------------------
PRONOUNS = {
    "ich", "du", "er", "sie", "es", "wir", "ihr", "Sie",
    "Ich", "Du", "Er", "Sie", "Es", "Wir", "Ihr",
}

def is_lexicon_candidate(text: str) -> bool:
    """
    True  => keep as lexicon entry (word/short phrase)
    False => route to examples.csv (sentences, questions, instructions)
    """
    t = (text or "").strip()
    if not t:
        return False

    # reject if it looks like a sentence/instruction
    if any(ch in t for ch in ["?", "!", ".", ":", ";"]):
        return False

    # reject obvious multi-clause/combined prompts
    if "," in t or " und " in t.lower() or " oder " in t.lower():
        return False

    tokens = t.split()

    # too long => probably sentence/instruction
    if len(tokens) > 4:
        return False

    # pronoun-led conjugation/sentences => examples, not lexicon
    if tokens and tokens[0] in PRONOUNS:
        return False

    return True


# -----------------------------
# Data model
# -----------------------------
@dataclass(frozen=True)
class ExtractedLine:
    ppt_file: str
    slide_number: int
    raw_text: str
    norm_text: str


# -----------------------------
# Extraction
# -----------------------------
def iter_pptx_files(root: Path) -> Iterable[Path]:
    yield from root.rglob("*.pptx")


def extract_lines_from_pptx(pptx_path: Path) -> List[ExtractedLine]:
    pres = Presentation(str(pptx_path))
    out: List[ExtractedLine] = []

    for i, slide in enumerate(pres.slides, start=1):
        texts: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)

        joined = "\n".join(texts)
        for raw_line in joined.splitlines():
            norm = normalize_line(raw_line)
            if norm:
                out.append(ExtractedLine(
                    ppt_file=pptx_path.name,
                    slide_number=i,
                    raw_text=raw_line,
                    norm_text=norm
                ))
    return out


def normalize_line(s: str) -> str:
    s = (s or "").strip()
    if not s:
        return ""

    # normalize bullets / weird chars, collapse whitespace
    s = s.replace("•", " ").replace("\u2022", " ").replace("\uf04c", " ")
    s = re.sub(r"\s+", " ", s).strip()

    # drop junk lines that are purely decorative
    if s in {"—", "-", "–"}:
        return ""

    return s


# -----------------------------
# Categorization (lexicon-oriented)
# -----------------------------
ARTICLES = {"der", "die", "das", "ein", "eine", "einen", "einem", "einer", "eines"}
TIME_WORDS = {
    "uhr", "heute", "morgen", "übermorgen", "gestern",
    "morgens", "abends", "nachts",
    "montag", "dienstag", "mittwoch", "donnerstag", "freitag", "samstag", "sonntag",
    "woche", "jahr", "monat", "tag", "stunden", "minute", "minuten",
    "vormittag", "mittag", "nachmittag", "abend", "nacht",
    "wochenende",
}

NUMBER_RE = re.compile(r"^\d+([:.,]\d+)?$")    # 12, 12:30, 12.5
TIME_RE = re.compile(r"^\d{1,2}:\d{2}$")       # 9:15

ADVERB_SET = {"nicht", "gern", "gerne", "sehr", "auch", "oft", "immer", "nie", "doch"}

def classify(line: str) -> str:
    """
    Returns one of:
      noun, verb, adjective, adverb, time_numbers, unknown
    NOTE: This is a *labeling heuristic*, not a grammar engine.
    """
    t = (line or "").strip()
    if not t:
        return "unknown"

    low = t.lower()
    parts = low.split()

    # numbers/time first
    if NUMBER_RE.match(low) or TIME_RE.match(low):
        return "time_numbers"

    if any(tok in TIME_WORDS for tok in parts):
        return "time_numbers"

    # adverbs (single tokens)
    if low in ADVERB_SET and " " not in low:
        return "adverb"

    # noun phrases: article + word (keep short phrases too)
    # e.g., "das Geld", "die Katze"
    if len(parts) >= 2 and parts[0] in ARTICLES:
        return "noun"

    # verbs: single infinitive candidates (gehen, kaufen, lernen)
    # very rough heuristic: ends with -en or -n and is one token
    if " " not in low and (low.endswith("en") or low.endswith("n")):
        return "verb"

    # adjectives: single token ending patterns (rough)
    if " " not in low and low.endswith(("ig", "lich", "isch")):
        return "adjective"

    return "unknown"


# -----------------------------
# Writers
# -----------------------------
def write_master_csv(lines: List[ExtractedLine], out_path: Path) -> None:
    out_path.parent.mkdir(parents=True, exist_ok=True)
    with out_path.open("w", newline="", encoding="utf-8") as f:
        w = csv.writer(f)
        w.writerow(["ppt_file", "slide_number", "norm_text", "category"])
        for L in lines:
            w.writerow([L.ppt_file, L.slide_number, L.norm_text, classify(L.norm_text)])


def write_examples_csv(rows: List[Tuple[str, int, str, str]], out_path: Path) -> None:
    out_path.parent.mkdir(parents=True, exist_ok=True)
    with out_path.open("w", newline="", encoding="utf-8") as f:
        w = csv.writer(f)
        w.writerow(["ppt_file", "slide_number", "norm_text", "category"])
        for ppt_file, slide, text, cat in rows:
            w.writerow([ppt_file, slide, text, cat])


def write_lexicon_yamls(lines: List[ExtractedLine], out_dir: Path) -> None:
    out_dir.mkdir(parents=True, exist_ok=True)

    # buckets store: text -> {"text":..., "sources":[...]} with deduped sources
    buckets: Dict[str, Dict[str, dict]] = {
        "noun": {},
        "verb": {},
        "adjective": {},
        "adverb": {},
        "time_numbers": {},
    }

    # (category, text) -> set of (ppt, slide) to dedupe sources
    seen_sources: Dict[Tuple[str, str], set] = {}

    for L in lines:
        text = L.norm_text
        cat = classify(text)

        # Only write lexicon YAMLs for lexicon candidates in these categories
        if cat not in buckets:
            continue
        if not is_lexicon_candidate(text):
            continue

        key = text
        if key not in buckets[cat]:
            buckets[cat][key] = {"text": text, "sources": []}

        ss_key = (cat, key)
        if ss_key not in seen_sources:
            seen_sources[ss_key] = set()

        src = (L.ppt_file, L.slide_number)
        if src not in seen_sources[ss_key]:
            buckets[cat][key]["sources"].append({"ppt": L.ppt_file, "slide": L.slide_number})
            seen_sources[ss_key].add(src)

    for cat, items_map in buckets.items():
        payload = {
            "schema_version": "0.1",
            "source": "pptx_extraction",
            "category": cat,
            "items": list(items_map.values()),
        }
        out_file = out_dir / f"{cat}s.yaml"  # nouns.yaml, verbs.yaml, adjectives.yaml, adverbs.yaml, time_numbers.yaml
        with out_file.open("w", encoding="utf-8") as f:
            yaml.safe_dump(payload, f, sort_keys=False, allow_unicode=True)


def build_examples(lines: List[ExtractedLine]) -> List[Tuple[str, int, str, str]]:
    """
    Anything NOT suitable for lexicon YAML goes here.
    Keeps category label (unknown/verb/noun/etc.) for later analysis.
    """
    ex: List[Tuple[str, int, str, str]] = []
    for L in lines:
        cat = classify(L.norm_text)
        if (cat in {"noun", "verb", "adjective", "adverb", "time_numbers"}) and is_lexicon_candidate(L.norm_text):
            continue
        ex.append((L.ppt_file, L.slide_number, L.norm_text, cat))
    return ex


# -----------------------------
# Main
# -----------------------------
def main():
    import argparse
    ap = argparse.ArgumentParser()
    ap.add_argument("ppt_root", help="Folder containing pptx files")
    args = ap.parse_args()

    ppt_root = Path(args.ppt_root).expanduser().resolve()
    all_lines: List[ExtractedLine] = []

    for pptx_path in iter_pptx_files(ppt_root):
        all_lines.extend(extract_lines_from_pptx(pptx_path))

    # Outputs
    write_master_csv(all_lines, Path("docs/ppt_extracted/items.csv"))
    write_lexicon_yamls(all_lines, Path("linguistic_rules/lexicon"))

    examples = build_examples(all_lines)
    write_examples_csv(examples, Path("docs/ppt_extracted/examples.csv"))

    print(f"Extracted lines: {len(all_lines)}")
    print("Wrote: docs/ppt_extracted/items.csv")
    print("Wrote: docs/ppt_extracted/examples.csv")
    print("Wrote: linguistic_rules/lexicon/* (lexicon candidates only)")


if __name__ == "__main__":
    main()