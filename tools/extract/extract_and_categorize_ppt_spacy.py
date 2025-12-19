#!/usr/bin/env python3
"""
Extract German vocabulary from a master PPTX using spaCy, and write repo-style items.csv.

Output schema (repo):
  ppt_file,slide_number,norm_text,category

Design goals:
- Extract *vocabulary* (tokens) not sentences/phrases.
- Also emit article+noun combos (e.g., "das Geld") as an extra vocab entry.
- Force common classroom edge cases:
  - "kein/keine/keinen/..." = determiner (never verb)
  - "ist/bin/bist/sind/..." = verb (sein)
  - question words: wer/was/wen pronoun; wo/woher/wohin/wann/warum/wie adverb
- Store verbs as *infinitive lemma* (norm_text = lemma).
- Store nouns/adjectives/adverbs etc. as normalized token (lowercased).
"""

from __future__ import annotations

import argparse
import csv
import re
import unicodedata as ud
from pathlib import Path
from typing import List, Optional, Tuple

import spacy
from pptx import Presentation


# =============================================================================
# spaCy model
# =============================================================================
# Install if needed:
#   python -m spacy download de_core_news_lg
# or (smaller):
#   python -m spacy download de_core_news_sm
DEFAULT_SPACY_MODEL = "de_core_news_lg"


# =============================================================================
# Robust classroom German overrides
# =============================================================================
NEG_DETERMINERS = {"kein", "keine", "keinen", "keinem", "keiner", "keines"}

ARTICLES = {
    "der", "die", "das",
    "ein", "eine",
    "den", "dem", "des",
    "einen", "einem", "einer", "eines",
}

# Auxiliaries / very common forms that spaCy can mis-tag in short examples
AUX_FORMS = {"bin", "bist", "ist", "sind", "seid", "war", "waren", "wäre", "wären"}
AUX_LEMMAS = {"sein", "haben", "werden"}

NEG_PARTICLES = {"nicht", "nie", "nichts", "kein"}  # 'kein' handled earlier anyway

QUESTION_ADVERBS = {"wo", "woher", "wohin", "wann", "warum", "wie"}
QUESTION_PRONOUNS = {"wer", "wen", "was"}

# Time/date helpers often used like adverbs / time-words
TIME_WORDS = {"halb", "viertel"}  # expand later if needed


# =============================================================================
# Verb classification (your detailed lists kept intact)
# =============================================================================
MODAL_VERBS = {"dürfen", "können", "mögen", "müssen", "sollen", "wollen"}

HIGHLY_IRREGULAR_VERBS = {
    "sein", "haben", "werden", "wissen", "tun", "gehen", "stehen", "bringen"
}

IRREGULAR_VERBS = {
    "bleiben", "leihen", "meiden", "preisen", "reiben", "scheinen",
    "schreiben", "schreien", "schweigen", "steigen", "treiben", "verzeihen",

    "beißen", "gleiten", "greifen", "kneifen", "pfeifen", "reißen",
    "reiten", "schleichen", "schmeißen", "schneiden", "streichen", "streiten",

    "biegen", "bieten", "fliegen", "fliehen", "fließen", "frieren",
    "gießen", "kriechen", "riechen", "schieben", "schießen", "schließen",
    "verlieren", "wiegen", "ziehen",

    "binden", "dringen", "finden", "gelingen", "klingen", "ringen",
    "singen", "sinken", "springen", "stinken", "trinken", "verschwinden",
    "zwingen",

    "befehlen", "bergen", "bersten", "brechen", "empfehlen", "gelten",
    "helfen", "nehmen", "schelten", "sprechen", "stechen", "sterben",
    "treffen", "verderben", "werben", "werfen",

    "essen", "fressen", "geben", "geschehen", "lesen", "messen",
    "sehen", "treten", "vergessen",

    "backen", "fahren", "graben", "laden", "schlagen", "tragen",
    "wachsen", "waschen",

    "blasen", "braten", "fallen", "fangen", "halten", "hängen",
    "lassen", "raten", "schlafen",

    "laufen", "saufen",

    "kommen",

    "lügen", "trügen",

    "rufen", "stoßen",

    "beginnen", "gewinnen", "liegen", "schwimmen", "sitzen", "spinnen",
}

STEM_CHANGING_VERBS = {
    "geben", "nehmen", "sprechen", "treffen", "helfen", "sterben",
    "werfen", "brechen", "erschrecken", "messen", "treten", "vergessen",
    "essen", "fressen", "gelten", "schelten", "stechen",

    "sehen", "lesen", "geschehen", "empfehlen", "befehlen", "stehlen",

    "fahren", "schlafen", "fallen", "fangen", "halten", "lassen",
    "raten", "schlagen", "tragen", "wachsen", "waschen", "laden",
    "graben", "braten", "backen",

    "laufen", "saufen",

    "stoßen",
}

SEPARABLE_PREFIXES = {
    "ab", "an", "auf", "aus", "bei", "ein", "fest", "her", "hin",
    "los", "mit", "nach", "vor", "weg", "weiter", "zu", "zurück",
    "zusammen", "durch", "über", "um", "unter", "wieder"
}

INSEPARABLE_PREFIXES = {"be", "emp", "ent", "er", "ge", "miss", "ver", "zer"}

MIXED_VERBS = {
    "brennen", "bringen", "denken", "kennen", "nennen", "rennen",
    "senden", "wenden", "wissen"
}


# =============================================================================
# Normalization helpers
# =============================================================================
WORDISH_RE = re.compile(r"^[A-Za-zÄÖÜäöüß]+(?:-[A-Za-zÄÖÜäöüß]+)*$")

def nfc_lower(s: str) -> str:
    return ud.normalize("NFC", s).strip().lower()

def is_vocab_token(text: str) -> bool:
    """
    Filter out junk:
    - punctuation
    - pure numbers
    - empty
    - tokens that are not "wordish" (German letters/hyphen)
    """
    t = nfc_lower(text)
    if not t:
        return False
    if t.isdigit():
        return False
    return bool(WORDISH_RE.match(t))


# =============================================================================
# Classroom meta / slide-instruction tokens to exclude
# =============================================================================
# These are not German vocab words; they come from grammar explanation slides.
META_TOKENS = {
    "infinitive", "conjugated", "conjugation",
    "direct", "object",
    "definite", "indefinite",
    "modal", "stem", "plural", "singular",
    "word", "order", "negation",
    "article", "articles",
    "pronoun", "preposition",
    "adjective", "adverb",
}

def is_meta_token_text(t_norm: str) -> bool:
    # Keep it simple and ruthless: if it’s one of these, it’s not vocab.
    return t_norm in META_TOKENS


# =============================================================================
# Verb helpers
# =============================================================================
def get_base_verb(lemma: str) -> Tuple[str, Optional[str], Optional[bool]]:
    """Return (base_verb, prefix, is_separable) if a prefix is detected."""
    lemma = nfc_lower(lemma)
    for prefix in SEPARABLE_PREFIXES:
        if lemma.startswith(prefix) and len(lemma) > len(prefix) + 2:
            return lemma[len(prefix):], prefix, True
    for prefix in INSEPARABLE_PREFIXES:
        if lemma.startswith(prefix) and len(lemma) > len(prefix) + 2:
            return lemma[len(prefix):], prefix, False
    return lemma, None, None

def classify_verb(lemma: str) -> str:
    lemma = nfc_lower(lemma)
    base, prefix, is_sep = get_base_verb(lemma)

    if base in MODAL_VERBS or lemma in MODAL_VERBS:
        verb_type = "modal verb"
    elif base in HIGHLY_IRREGULAR_VERBS or lemma in HIGHLY_IRREGULAR_VERBS:
        verb_type = "highly irregular verb"
    elif base in MIXED_VERBS or lemma in MIXED_VERBS:
        verb_type = "mixed verb"
    elif base in IRREGULAR_VERBS or lemma in IRREGULAR_VERBS:
        if base in STEM_CHANGING_VERBS or lemma in STEM_CHANGING_VERBS:
            verb_type = "irregular stem-changing verb"
        else:
            verb_type = "irregular verb"
    elif base in STEM_CHANGING_VERBS or lemma in STEM_CHANGING_VERBS:
        verb_type = "stem-changing verb"
    else:
        verb_type = "regular verb"

    if prefix:
        if is_sep:
            return f"{verb_type} (separable: {prefix}-)"
        return f"{verb_type} (inseparable: {prefix}-)"
    return verb_type


# =============================================================================
# Category mapping: spaCy -> repo buckets
# =============================================================================
def get_internal_category(token) -> str:
    """
    Internal categories (fine-grained) then mapped down to repo categories.
    """
    t = nfc_lower(token.text)
    lem = nfc_lower(token.lemma_)

    # HARD OVERRIDES FIRST (prevents stupid mistakes like kein=verb)
    if t in NEG_DETERMINERS:
        return "determiner_negation"
    if t in ARTICLES:
        return "article"
    if t in QUESTION_PRONOUNS:
        return "pronoun"
    if t in QUESTION_ADVERBS:
        return "adverb"
    if t in NEG_PARTICLES:
        return "particle"
    if t in TIME_WORDS:
        return "time_word"

    # Auxiliary forms/lemmas
    if t in AUX_FORMS or lem in AUX_LEMMAS:
        return "verb_aux"

    # spaCy POS-based fallback
    pos = token.pos_
    if pos == "DET":
        return "article"
    if pos == "NOUN":
        return "noun"
    if pos == "PROPN":
        return "proper_noun"
    if pos in {"VERB", "AUX"}:
        return classify_verb(lem)
    if pos == "ADJ":
        return "adjective"
    if pos == "ADV":
        return "adverb"
    if pos == "PRON":
        return "pronoun"
    if pos == "ADP":
        return "preposition"
    if pos in {"CCONJ", "CONJ"}:
        return "conjunction"
    if pos == "SCONJ":
        return "subordinating_conjunction"
    if pos == "NUM":
        return "numeral"
    if pos == "INTJ":
        return "interjection"

    return "other"


def map_internal_to_repo_category(internal_cat: str) -> str:
    """
    Map internal categories to repo categories.
    Repo buckets observed/expected:
      noun, verb, adjective, time_numbers, unknown
    """
    ic = internal_cat

    if ic in {"noun", "proper_noun"}:
        return "noun"

    if ic == "adjective":
        return "adjective"

    if ic == "time_word" or ic == "numeral":
        return "time_numbers"

    # determiners/articles/particles/pronouns/etc → unknown (unless you add buckets later)
    if ic in {
        "article", "determiner_negation",
        "adverb", "preposition", "particle", "pronoun",
        "conjunction", "subordinating_conjunction",
        "interjection", "other",
    }:
        return "unknown"

    # Any verb classification string contains "verb"
    if "verb" in ic:
        return "verb"

    return "unknown"


# =============================================================================
# PPTX extraction
# =============================================================================
def iter_slide_texts(pptx_path: Path) -> List[Tuple[int, str]]:
    prs = Presentation(str(pptx_path))
    out: List[Tuple[int, str]] = []

    for slide_num, slide in enumerate(prs.slides, 1):
        chunks: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                chunks.append(shape.text)
        txt = " ".join(chunks)
        out.append((slide_num, txt))
    return out


def extract_vocab_rows(nlp, pptx_path: Path) -> List[dict]:
    ppt_name = pptx_path.name
    rows: List[dict] = []

    slide_texts = iter_slide_texts(pptx_path)

    for slide_num, full_text in slide_texts:
        if not full_text.strip():
            continue

        doc = nlp(full_text)

        prev_token = None  # for article+noun combos
        prev_internal = None

        for tok in doc:
            if tok.is_space or tok.is_punct:
                continue

            raw = tok.text
            if not is_vocab_token(raw):
                prev_token = None
                prev_internal = None
                continue

            t_norm = nfc_lower(raw)
            if is_meta_token_text(t_norm):
                # This is the main fix for your grep “infinitive” pollution.
                prev_token = None
                prev_internal = None
                continue

            lem_norm = nfc_lower(tok.lemma_)
            internal_cat = get_internal_category(tok)

            # Decide what we store as norm_text:
            # - verbs: store infinitive lemma
            # - others: store normalized token text
            is_verbish = "verb" in internal_cat
            store_text = lem_norm if is_verbish else t_norm

            # Also protect against meta tokens showing up as lemmas
            if is_meta_token_text(store_text):
                prev_token = None
                prev_internal = None
                continue

            repo_cat = map_internal_to_repo_category(internal_cat)

            # Emit token itself
            rows.append({
                "ppt_file": ppt_name,
                "slide_number": slide_num,
                "norm_text": store_text,
                "category": repo_cat,
            })

            # Emit article+noun combo (extra vocab entry)
            # Only when we see article/determiner_negation followed immediately by noun/proper noun.
            if prev_token is not None and prev_internal in {"article", "determiner_negation"}:
                if internal_cat in {"noun", "proper_noun"}:
                    combo = f"{nfc_lower(prev_token.text)} {t_norm}"
                    if not is_meta_token_text(combo):
                        rows.append({
                            "ppt_file": ppt_name,
                            "slide_number": slide_num,
                            "norm_text": combo,
                            "category": "noun",
                        })

            prev_token = tok
            prev_internal = internal_cat

    return rows


def write_items_csv(rows: List[dict], out_csv: Path) -> None:
    out_csv.parent.mkdir(parents=True, exist_ok=True)

    with out_csv.open("w", encoding="utf-8", newline="") as f:
        w = csv.DictWriter(f, fieldnames=["ppt_file", "slide_number", "norm_text", "category"])
        w.writeheader()
        for r in rows:
            w.writerow(r)


# =============================================================================
# CLI
# =============================================================================
def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--pptx", required=True, help="Path to master PPTX file (all slides merged).")
    ap.add_argument("--out", default="docs/ppt_extracted/items.csv", help="Output items.csv path.")
    ap.add_argument("--model", default=DEFAULT_SPACY_MODEL, help="spaCy German model name.")
    args = ap.parse_args()

    pptx_path = Path(args.pptx)
    if not pptx_path.exists():
        raise SystemExit(f"ERROR: PPTX not found: {pptx_path}")

    out_csv = Path(args.out)

    try:
        nlp = spacy.load(args.model)
    except OSError as e:
        raise SystemExit(
            f"ERROR: spaCy model '{args.model}' not found.\n"
            f"Install with:\n"
            f"  python -m spacy download {args.model}\n"
            f"(or switch to de_core_news_sm)\n"
        ) from e

    print("Parsing PowerPoint presentation...")
    rows = extract_vocab_rows(nlp, pptx_path)

    print(f"Extracted {len(rows)} vocab rows (includes article+noun combos).")
    write_items_csv(rows, out_csv)
    print(f"Wrote: {out_csv}")

    print("\nNext steps (repo pipeline):")
    print("  python tools/rebuild_lexicon_from_items_csv.py")
    print("  python tools/add_lex_ids_and_build_index.py")
    print("  python tools/link_conjugation_tables_to_lexicon.py")
    print("  python tools/validate_ruleset_integrity.py")


if __name__ == "__main__":
    main()
