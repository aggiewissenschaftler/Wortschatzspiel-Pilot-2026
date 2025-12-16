from pathlib import Path
from pptx import Presentation

PPT_DIR = Path("source_material/powerpoints")
OUT = Path("docs/ppt_text_dump.md")
OUT.parent.mkdir(parents=True, exist_ok=True)

def extract_text(ppt_path: Path) -> list[str]:
    prs = Presentation(str(ppt_path))
    lines = [f"# {ppt_path.name}"]
    for i, slide in enumerate(prs.slides, start=1):
        slide_lines = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_lines.append(shape.text.strip())
        if slide_lines:
            lines.append(f"\n## Slide {i}\n" + "\n\n".join(slide_lines))
    return lines

def main():
    all_lines = []
    for ppt in sorted(PPT_DIR.glob("*.pptx")):
        all_lines.extend(extract_text(ppt))
        all_lines.append("\n---\n")
    OUT.write_text("\n".join(all_lines), encoding="utf-8")
    print(f"Wrote: {OUT}")

if __name__ == "__main__":
    main()